#!/usr/bin/env python3
"""Generate the 3-minute, 6-slide talk deck (DeepCodonOpt_slides.pptx).

Self-contained: only needs python-pptx.
    pip install python-pptx
    python build_slides.py            # -> DeepCodonOpt_slides.pptx

Matches paper/main.tex and paper/slides_content.md (numbers, colours, speaker notes).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (matches the paper figures) ----
NAVY   = RGBColor(0x21, 0x29, 0x39)
PURPLE = RGBColor(0x81, 0x72, 0xB3)   # Method B
RED    = RGBColor(0xC4, 0x4E, 0x52)   # baseline / Method A
BLUE   = RGBColor(0x4C, 0x72, 0xB0)   # natural
GRAY   = RGBColor(0x55, 0x5B, 0x68)
LIGHT  = RGBColor(0xF1, 0xF2, 0xF6)
LILAC  = RGBColor(0xEC, 0xE8, 0xF6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _runs_font(p, size, bold, color, name="Calibri"):
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = name


def rect(slide, l, t, w, h, color, line=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if not line:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tf = slide.shapes.add_textbox(l, t, w, h).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def header(slide, kicker, title, kcolor=PURPLE):
    rect(slide, 0, 0, Inches(0.28), SH, kcolor)            # left accent bar
    tf = textbox(slide, Inches(0.7), Inches(0.45), Inches(12.0), Inches(1.4))
    p = tf.paragraphs[0]; p.text = kicker
    _runs_font(p, 14, True, kcolor)
    p2 = tf.add_paragraph(); p2.text = title; p2.space_before = Pt(2)
    _runs_font(p2, 30, True, NAVY)


def bullets(slide, items, l, t, w, h, size=19, gap=10):
    tf = textbox(slide, l, t, w, h)
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if lvl == 0 else "–  ") + text
        p.level = lvl
        p.space_after = Pt(gap)
        _runs_font(p, size if lvl == 0 else size - 2, False, NAVY if lvl == 0 else GRAY)
        if lvl == 1:
            p.runs[0].font.color.rgb = GRAY
    return tf


def table(slide, data, l, t, w, h, highlight=(), hl_color=PURPLE, header_fill=NAVY):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, l, t, w, h).table
    gt.first_row = True
    for ri, row in enumerate(data):
        gt.rows[ri].height = Inches(h.inches / rows)
        for ci, val in enumerate(row):
            c = gt.cell(ri, ci)
            c.text = str(val)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(3); c.margin_bottom = Pt(3)
            c.margin_left = Pt(8); c.margin_right = Pt(8)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()   # empty cells have no run
            run.font.name = "Calibri"; run.font.size = Pt(15)
            if ri == 0:
                run.font.bold = True; run.font.color.rgb = WHITE
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT if ri % 2 == 0 else WHITE
                run.font.color.rgb = NAVY
                if ri in highlight:
                    run.font.bold = True; run.font.color.rgb = hl_color
                    c.fill.fore_color.rgb = LILAC
    return gt


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def callout(slide, text, l, t, w, h, fill=LILAC, color=PURPLE):
    box = rect(slide, l, t, w, h, fill)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(14); tf.margin_right = Pt(14)
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    _runs_font(p, 18, True, color)


# =========================================================================== Slide 1
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.06), PURPLE)
tf = textbox(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(3.0))
p = tf.paragraphs[0]; p.text = "DeepCodonOpt"
_runs_font(p, 52, True, WHITE)
p2 = tf.add_paragraph(); p2.text = "Neural Codon Optimization with Biological Constraints"
p2.space_before = Pt(10); _runs_font(p2, 26, False, RGBColor(0xC8, 0xCE, 0xDC))
tf2 = textbox(s, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.5))
p = tf2.paragraphs[0]
p.text = "Training-time constraint loss  vs.  inference-time constrained decoding"
_runs_font(p, 18, True, PURPLE)
p = tf2.add_paragraph(); p.text = "Isu Kim · San Kim · Hoseok Lee · Daehee Hong"
p.space_before = Pt(10); _runs_font(p, 15, False, RGBColor(0x9A, 0xA2, 0xB4))
notes(s, "Neural codon optimizers write great DNA for expression, but ignore whether it can "
         "actually be manufactured. We ask where to put the manufacturing rules: during training "
         "or at inference.")

# =========================================================================== Slide 2
s = prs.slides.add_slide(BLANK)
header(s, "DIAGNOSIS", "Great codons, terrible to synthesize")
bullets(s, [
    ("Scored 3,000 E. coli genes with the IDT complexity score (synthesis-difficulty index)", 0),
    ("Baseline model: CAI 0.89 > natural 0.74  —  codon usage is excellent…", 0),
    ("…but IDT complexity 4.9 vs 1.9  —  much harder to synthesize", 0),
    ("Cause is one issue: Overall Repeat  24% → 61%", 0),
], Inches(0.7), Inches(1.9), Inches(7.1), Inches(4.6), size=19)
table(s, [
    ["", "complexity", "CAI", "Overall Rep."],
    ["Natural", "1.91", "0.735", "23.7%"],
    ["Baseline (FT)", "4.91", "0.888", "60.8%"],
    ["Pretrained", "6.63", "0.885", "63.9%"],
], Inches(8.1), Inches(2.5), Inches(4.6), Inches(2.4), highlight=(2,), hl_color=RED)
callout(s, "Greedy decoding repeats each amino acid's favorite codon → repetitive DNA",
        Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.75))
notes(s, "Greedy decoding always picks each amino acid's favorite codon — maximizing CAI but "
         "producing repetitive DNA that IDT flags as hard to synthesize.")

# =========================================================================== Slide 3
s = prs.slides.add_slide(BLANK)
header(s, "METHOD A — TRAINING-TIME", "Add a constraint loss → it backfires", kcolor=RED)
bullets(s, [
    ("Loss = MLM + λ · (repeat/GC penalty on the model's soft distribution)", 0),
    ("Result: complexity 4.9 → 9.3,  CAI 0.89 → 0.43", 0),
    ("Breakdown: GC issues → 0%, but Overall Repeat 61% → 80%", 0),
    ("It traded one issue for a much worse one", 1),
], Inches(0.7), Inches(1.9), Inches(7.1), Inches(4.6), size=19)
table(s, [
    ["", "complexity", "CAI", "Overall Rep."],
    ["Baseline", "4.91", "0.888", "60.8%"],
    ["Method A", "9.32", "0.426", "79.7%"],
], Inches(8.1), Inches(2.6), Inches(4.6), Inches(1.9), highlight=(2,), hl_color=RED)
callout(s, "Every λ ∈ {1,3,10,30} gave the same result", Inches(8.1), Inches(4.9),
        Inches(4.6), Inches(0.75), fill=RGBColor(0xFA, 0xEC, 0xEC), color=RED)
notes(s, "Adding a differentiable constraint loss made everything worse. Why? The penalty scores "
         "the model's probabilities, not the sequence it actually writes.")

# =========================================================================== Slide 4
s = prs.slides.add_slide(BLANK)
header(s, "METHOD A — WHY IT FAILS", "The loss measures the wrong thing", kcolor=RED)
bullets(s, [
    ("The penalty acts on the probability distribution, not the emitted sequence", 0),
    ("The model satisfies it by flattening its distribution (→ rare codons)…", 0),
    ("…at ≈ zero MLM cost, so λ doesn't matter", 1),
    ("But the decoded DNA is unchanged or worse — surrogate ≠ objective", 0),
], Inches(0.7), Inches(2.0), Inches(11.9), Inches(3.0), size=20)
callout(s, "Like lowering your test average by leaving answers blank, "
           "not by getting them right.", Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.0))
notes(s, "It's like lowering your test average by leaving answers blank instead of getting them "
         "right. The soft training objective and the real objective are misaligned, so minimizing "
         "the loss moves away from the goal.")

# =========================================================================== Slide 5
s = prs.slides.add_slide(BLANK)
header(s, "METHOD B — INFERENCE-TIME", "Constrained temperature decoding → it works")
bullets(s, [
    ("Temperature  P(c) ∝ exp(z/T):  diversifies codons → breaks repeats", 0),
    ("Masking: delete codons that form restriction sites / homopolymers", 0),
    ("Sweep T → clean Pareto; knee at T ≈ 0.7–0.8", 0),
    ("beats natural on BOTH axes (complexity 1.6 < 1.9, CAI 0.78 > 0.74)", 1),
], Inches(0.7), Inches(1.9), Inches(7.0), Inches(4.6), size=19)
table(s, [
    ["T", "CAI", "complexity", "Overall Rep."],
    ["natural", "0.735", "1.91", "23.7%"],
    ["0.5", "0.812", "2.20", "32.6%"],
    ["0.7", "0.781", "1.59", "23.9%"],
    ["0.8", "0.767", "1.41", "21.3%"],
    ["1.0", "0.748", "1.18", "17.4%"],
], Inches(8.0), Inches(1.95), Inches(4.7), Inches(3.3), highlight=(3, 4), hl_color=PURPLE)
callout(s, "Temperature is the driver; masking is the guarantee",
        Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7))
notes(s, "Change only the decoding: sample at temperature to break repeats, and mask out illegal "
         "codons. This operates on the real sequence, so it actually reduces synthesis complexity "
         "— and at T around 0.7 to 0.8 it beats natural genes on both axes.")

# =========================================================================== Slide 6
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.95), SW, Inches(0.06), PURPLE)
tf = textbox(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(2.0))
p = tf.paragraphs[0]; p.text = "Takeaway"
_runs_font(p, 22, True, PURPLE)
p = tf.add_paragraph()
p.text = "Discrete constraints belong on the discrete output, at decode time."
p.space_before = Pt(8); _runs_font(p, 30, True, WHITE)
bullets_tf = textbox(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(3.0))
for i, t in enumerate([
    "Training-time soft loss → gameable, backfires (4.9 → 9.3)",
    "Inference-time masking + temperature → exact + effective, surpasses natural",
    "Future: look-ahead decoding for long-range repeats",
]):
    p = bullets_tf.paragraphs[0] if i == 0 else bullets_tf.add_paragraph()
    p.text = "•  " + t; p.space_after = Pt(12)
    _runs_font(p, 19, False, RGBColor(0xD6, 0xDB, 0xE6))
tf2 = textbox(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8))
p = tf2.paragraphs[0]
p.text = "github.com/<user>/DeepCodonOpt    ·    Zenodo 10.5281/zenodo.20602294"
_runs_font(p, 15, False, RGBColor(0x9A, 0xA2, 0xB4))
notes(s, "One principle: don't relax discrete biological constraints into a training loss — "
         "enforce them on the actual sequence at decode time.")

# ---------------------------------------------------------------------------
out = "DeepCodonOpt_slides.pptx"
prs.save(out)
print("wrote", out, "—", len(prs.slides), "slides")
